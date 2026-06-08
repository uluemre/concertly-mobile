"""
Concertly — 12 Slaytlık Yatırımcı / Jüri Pitch Deck
Blueprint: CONCERTLY_PRESENTATION_BLUEPRINT.md
Tasarım dili: Apple Keynote sadeliği + Airbnb duygusallığı + Spotify Wrapped renk enerjisi
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Renk Sistemi ──────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x0D, 0x1A)   # Ana arka plan
SURFACE = RGBColor(0x1A, 0x1A, 0x2E)   # Kart zemini
SURFACE2= RGBColor(0x16, 0x16, 0x28)   # İkincil kart
BORDER  = RGBColor(0x2A, 0x2A, 0x3E)   # Kart çerçevesi
SHADOW  = RGBColor(0x07, 0x07, 0x10)   # Gölge

RED     = RGBColor(0xE9, 0x45, 0x60)   # Ana vurgu — kırmızı-pembe
TEAL    = RGBColor(0x00, 0xD4, 0xAA)   # İkincil vurgu — teal
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)   # Üçüncül — mor
ORANGE  = RGBColor(0xF5, 0xA6, 0x23)   # Turuncu

WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x8B, 0x8B, 0xA0)   # Soluk metin
DARK    = RGBColor(0x22, 0x22, 0x38)   # Çok koyu arka plan unsuru

# ── Slayt Boyutu ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = 13.33, 7.5

# ── Ekran Yolları ─────────────────────────────────────────────────────────────
_S = r"C:\Users\EMRE\Desktop\concertly-mobile\screenshots"
SCR = {
    "login":       _S + r"\WhatsApp Image 2026-06-04 at 11.15.41.jpeg",
    "home":        _S + r"\WhatsApp Image 2026-06-04 at 11.04.56.jpeg",
    "badges":      _S + r"\WhatsApp Image 2026-06-04 at 11.04.56 (1).jpeg",
    "passport":    _S + r"\WhatsApp Image 2026-06-04 at 11.04.56 (2).jpeg",
    "music":       _S + r"\WhatsApp Image 2026-06-04 at 11.04.56 (3).jpeg",
    "communities": _S + r"\WhatsApp Image 2026-06-04 at 11.04.57.jpeg",
    "buddy":       _S + r"\WhatsApp Image 2026-06-04 at 11.04.57 (2).jpeg",
    "map_tr":      _S + r"\WhatsApp Image 2026-06-04 at 11.04.57 (3).jpeg",
    "events":      _S + r"\WhatsApp Image 2026-06-04 at 11.04.58 (1).jpeg",
    "feed":        _S + r"\WhatsApp Image 2026-06-04 at 11.04.58 (2).jpeg",
}

# ── Yardımcı Fonksiyonlar ─────────────────────────────────────────────────────

def apply_grad(shape, stops, ang=5400000):
    """Shape'e gradient fill uygular."""
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp, qn('p:spPr'))
    for old in spPr.findall(qn('a:solidFill')) + spPr.findall(qn('a:gradFill')):
        spPr.remove(old)
    gf  = etree.SubElement(spPr, qn('a:gradFill'))
    gsL = etree.SubElement(gf, qn('a:gsLst'))
    for pos, rgb in stops:
        gs  = etree.SubElement(gsL, qn('a:gs'), pos=str(int(pos * 100000)))
        clr = etree.SubElement(gs, qn('a:srgbClr'),
                               val=f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")
    etree.SubElement(gf, qn('a:lin'), ang=str(ang), scaled='0')

def rgb_tuple(rgb_color):
    """RGBColor → (r,g,b) tuple"""
    h = str(rgb_color)
    return (int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def bg(slide):
    """Tam arka plan dikdörtgeni."""
    sh = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = BG
    sh.line.fill.background()
    return sh

def rect(slide, l, t, w, h, color=SURFACE, rounding=None):
    """Basit dikdörtgen."""
    shape_type = 5 if rounding else 1
    sh = slide.shapes.add_shape(
        shape_type, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    if rounding:
        sh.adjustments[0] = rounding
    return sh

def oval(slide, l, t, d, color=PURPLE):
    sh = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh

def card(slide, l, t, w, h, border_color=None, surface=SURFACE):
    """Kart: gölge + yüzey + opsiyonel border."""
    sd = slide.shapes.add_shape(
        1, Inches(l+0.06), Inches(t+0.07), Inches(w), Inches(h))
    sd.fill.solid(); sd.fill.fore_color.rgb = SHADOW
    sd.line.fill.background()
    sh = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = surface
    if border_color:
        sh.line.color.rgb = border_color
        sh.line.width = Inches(0.016)
    else:
        sh.line.fill.background()
    return sh

def txt(slide, text, l, t, w, h, size, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    b = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    b.word_wrap = wrap
    tf = b.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return b

def phone(slide, img_key, l, t, w):
    """iPhone mockup çerçevesiyle ekran görüntüsü."""
    ratio = 2048 / 945   # iPhone 15 Pro oranı
    h = w * ratio
    # Gölge
    sd = slide.shapes.add_shape(
        1, Inches(l+0.09), Inches(t+0.1), Inches(w+0.04), Inches(h+0.04))
    sd.fill.solid(); sd.fill.fore_color.rgb = SHADOW; sd.line.fill.background()
    # Çerçeve (biraz daha büyük)
    fr = slide.shapes.add_shape(
        5, Inches(l-0.07), Inches(t-0.07), Inches(w+0.14), Inches(h+0.14))
    fr.fill.solid(); fr.fill.fore_color.rgb = RGBColor(0x1C, 0x1C, 0x30)
    fr.line.color.rgb = BORDER; fr.line.width = Inches(0.018)
    # Görüntü
    img_path = SCR.get(img_key, "")
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        fb = slide.shapes.add_shape(
            1, Inches(l), Inches(t), Inches(w), Inches(h))
        fb.fill.solid(); fb.fill.fore_color.rgb = SURFACE2; fb.line.fill.background()
        txt(slide, img_key, l, t+h/2-0.2, w, 0.4, 11, color=MUTED, align=PP_ALIGN.CENTER)
    return h

def accent_line(slide, l, t, w=1.2, color=RED, thick=0.07):
    sh = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(thick))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def glow(slide, cx, cy, size, color, opacity_hex="12"):
    """Hafif radial parıltı efekti — büyük oval, çok düşük opacity."""
    sh = slide.shapes.add_shape(
        9, Inches(cx-size/2), Inches(cy-size/2), Inches(size), Inches(size))
    # Solid fill + transparency
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    # Şeffaflık için XML manipülasyonu
    sp  = sh._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp, qn('p:spPr'))
    fld  = spPr.find(qn('a:solidFill'))
    if fld is not None:
        srgb = fld.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn('a:alpha'))
            # pptx'te 100000 = %100, opacity_hex "12" ≈ %7
            pct = int(opacity_hex, 16) * 100000 // 255
            alpha.set('val', str(pct))
    sh.line.fill.background()

def pill_badge(slide, l, t, w, h, bg_color, text_str, tsize=10, tc=WHITE):
    sh = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = bg_color; sh.line.fill.background()
    txt(slide, text_str, l, t+0.02, w, h-0.04, tsize,
        bold=True, color=tc, align=PP_ALIGN.CENTER)

def gradient_rect(slide, l, t, w, h, c1, c2, angle=0):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c1; sh.line.fill.background()
    apply_grad(sh, [(0, rgb_tuple(c1)), (1, rgb_tuple(c2))], ang=angle)
    return sh

def slide_number_dot(slide, num, total=12):
    """Sağ alt köşede küçük slayt numarası."""
    txt(slide, f"{num} / {total}", W-1.2, H-0.38, 1.0, 0.3,
        9, color=MUTED, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 1 — KAPAK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# Arka plan glow efektleri
glow(s, 3.5, 4.5, 7.0,  PURPLE, "0E")
glow(s, 10.5, 2.0, 5.0, RED,    "0A")

# Sol dikey şerit
gradient_rect(s, 0, 0, 0.18, H, RED, PURPLE, angle=10800000)

# Logo ikonunu dene, yoksa emoji
icon_path = r"C:\Users\EMRE\Desktop\concertly-mobile\mobile\assets\icon.png"
if os.path.exists(icon_path):
    s.shapes.add_picture(icon_path, Inches(1.1), Inches(1.6), Inches(1.2), Inches(1.2))
else:
    oval(s, 1.1, 1.6, 1.2, PURPLE)
    txt(s, "🎵", 1.1, 1.62, 1.2, 1.2, 30, align=PP_ALIGN.CENTER)

# Ana başlık
txt(s, "CONCERTLY", 2.7, 1.58, 8.5, 1.3, 76, bold=True, color=WHITE)

# Gradient accent çizgi (kırmızı → mor)
gradient_rect(s, 2.7, 2.98, 5.8, 0.07, RED, PURPLE, angle=0)

# Alt başlık
txt(s, "Müziği birlikte yaşa.", 2.7, 3.14, 9.0, 0.6, 28, italic=True, color=MUTED)

# Etiketler
for i, (label, col) in enumerate([
    ("Keşfet", RED), ("·", MUTED), ("Birlikte Yaşa", TEAL), ("·", MUTED), ("Biriktir", PURPLE)
]):
    txt(s, label, 2.7 + i*1.82, 3.88, 1.7, 0.45, 16, bold=(label != "·"),
        color=col)

# Sağ taraf: Login ekranı mockup
phone(s, "login", 9.8, 0.55, 2.6)

# Alt bilgi
txt(s, "Emre  ·  2026", 2.7, 6.95, 5.0, 0.38, 11, color=MUTED)
slide_number_dot(s, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 2 — İNSAN GEREKÇESİ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# Feed ekranındaki konser fotoğrafını tam ekran arka plan gibi kullan
# Sol yarı koyu overlay + sağ yarı mockup
rect(s, 0, 0, W, H, BG)

# Sol: büyük duygusal metin
glow(s, 4.0, 3.5, 8.0, RED, "0C")

txt(s, "İnsanlar neden", 1.0, 1.3, 8.0, 0.75, 32, color=MUTED)
txt(s, "konsere gider?", 1.0, 2.05, 8.5, 1.1, 54, bold=True, color=WHITE)

accent_line(s, 1.0, 3.3, 2.0, RED, 0.07)

txt(s, "Müzik için değil.", 1.0, 3.55, 8.5, 0.75, 34, bold=True, color=WHITE)
txt(s, "O an için.",        1.0, 4.3,  8.5, 0.75, 34, bold=True, color=RED)

txt(s,
    "Binlerce kişiyle aynı anda, aynı ritme teslim olmak —\n"
    "bu deneyim çoğaltılamaz, ikinci kez tam olarak yaşanamaz.",
    1.0, 5.22, 7.5, 1.0, 16, color=MUTED, italic=True)

# Sağ: Feed ekranından konser fotoğrafı
phone(s, "feed", 9.65, 0.6, 2.7)

slide_number_dot(s, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 3 — PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
glow(s, 7.0, 5.5, 6.0, PURPLE, "0A")

txt(s, "Bugün bir konser planlamak için...", 1.0, 0.55, 11.0, 0.6, 22, color=MUTED)

# 5 adım akışı
steps = [
    ("🎟", "Biletix",   "Bilet ara"),
    ("📸", "Instagram", "Araştır"),
    ("💬", "WhatsApp",  "Biri bul"),
    ("📅", "Takvim",    "Not al"),
    ("⏰", "24 saat",   "Kaybol"),
]
step_w, step_h = 2.1, 2.2
total_flow = len(steps)*step_w + (len(steps)-1)*0.55
x0 = (W - total_flow) / 2

for i, (icon, label, sub) in enumerate(steps):
    lx = x0 + i*(step_w+0.55)
    is_last = (i == len(steps)-1)
    border_col = RED if is_last else BORDER
    srf = RGBColor(0x1F, 0x08, 0x12) if is_last else SURFACE
    card(s, lx, 1.35, step_w, step_h, border_col, srf)
    txt(s, icon,  lx,      1.52, step_w, 0.7,  32, align=PP_ALIGN.CENTER)
    txt(s, label, lx,      2.3,  step_w, 0.42, 14, bold=True,
        color=RED if is_last else WHITE, align=PP_ALIGN.CENTER)
    txt(s, sub,   lx,      2.75, step_w, 0.35, 11,
        color=MUTED, align=PP_ALIGN.CENTER)
    if not is_last:
        txt(s, "→", lx+step_w+0.08, 1.9, 0.45, 0.55, 24, bold=True, color=BORDER)

# Alt büyük mesaj
rect(s, 1.0, 3.88, W-2.0, 0.005, BORDER)

txt(s, "5 platform",  2.0,  4.15, 3.2, 0.75, 46, bold=True, color=WHITE)
txt(s, "1 deneyim",   5.4,  4.15, 3.2, 0.75, 46, bold=True, color=WHITE)
txt(s, "0 iz.",       8.8,  4.15, 2.8, 0.75, 46, bold=True, color=RED)

# Üç problem kart
probs = [
    ("📍", "Keşif dağınık",       "Kişiselleştirilmiş\netkinlik önerisi yok",   PURPLE),
    ("👤", "Yalnız gitme korkusu", "Konsere birlikte\ngidecek biri bulunamıyor", RED),
    ("🗂",  "Kayıp geçmiş",        "Konser anıları\nhiçbir yerde birikmez",     TEAL),
]
pcw = 3.7
px0 = (W - 3*pcw - 2*0.35) / 2
for i, (ico, t_str, d_str, col) in enumerate(probs):
    lx = px0 + i*(pcw+0.35)
    card(s, lx, 5.15, pcw, 1.95, col)
    txt(s, f"{ico}  {t_str}", lx+0.25, 5.35, pcw-0.3, 0.42, 13, bold=True, color=col)
    txt(s, d_str, lx+0.25, 5.82, pcw-0.3, 0.85, 12, color=MUTED)

slide_number_dot(s, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 4 — VİZYON
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# Arka plan glowlar
glow(s, 3.0,  4.0, 9.0, RED,    "0B")
glow(s, 11.0, 2.0, 7.0, PURPLE, "0A")
glow(s, 7.0,  7.0, 6.0, TEAL,   "08")

# Üç küçük etiket
for i, (word, col) in enumerate([("Keşfet.", RED), ("Birlikte yaşa.", TEAL), ("Biriktir.", PURPLE)]):
    txt(s, word, 1.5 + i*3.65, 1.55, 3.5, 0.6, 24, bold=True, color=col)

# Ortada ince çizgi
rect(s, 1.5, 2.32, W-3.0, 0.014, BORDER)

# Ana vizyon cümlesi
txt(s, "Canlı müzik deneyiminin", 1.5, 2.55, W-3.0, 1.1, 58, bold=True, color=WHITE)

# "dijital katmanı." — gradient efekt için iki ayrı txt
gradient_rect(s, 1.5, 3.65, 8.2, 0.92, RED, PURPLE, angle=0)
# Gradient dikdörtgen üzerine metin
txt(s, "dijital katmanı.", 1.5, 3.65, 8.2, 0.92, 58, bold=True,
    color=WHITE, align=PP_ALIGN.LEFT)

# Alt metin
txt(s,
    "Bir bilet platformu değil. Bir sosyal ağ değil.\n"
    "Canlı müzik deneyiminin tüm döngüsü — önce, sırasında, sonra.",
    1.5, 4.85, W-3.0, 0.9, 18, color=MUTED, italic=True)

slide_number_dot(s, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 5 — ÜRÜN: ANA EKRAN
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
glow(s, 3.0, 3.5, 7.0, PURPLE, "0C")

# Sol metin bloğu
txt(s, "Şehrine göre", 0.9, 1.1, 6.5, 0.72, 42, bold=True, color=WHITE)
txt(s, "etkinlikler.",  0.9, 1.83, 6.5, 0.72, 42, bold=True, color=RED)

accent_line(s, 0.9, 2.72, 1.5, RED, 0.065)

desc_items = [
    ("📍", "11 şehir · tür filtresi · anlık arama"),
    ("👥", "Arkadaşlarının planlarını gör"),
    ("⭐", "Öne çıkan etkinlikler carousel'i"),
    ("🎟", "Bilet linkine tek tıkla ulaş"),
]
for i, (ico, line) in enumerate(desc_items):
    txt(s, f"{ico}  {line}", 0.9, 2.98+i*0.58, 6.5, 0.52, 15, color=WHITE)

# İstatistik badge
card(s, 0.9, 5.68, 6.0, 0.95, TEAL, SURFACE2)
txt(s, "99 etkinlik  ·  21 post  ·  Türkiye geneli",
    1.1, 5.88, 5.8, 0.45, 17, bold=True, color=TEAL)

# Sağ: Home screen mockup
phone(s, "home", 8.35, 0.3, 3.8)

slide_number_dot(s, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 6 — CONCERT BUDDY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
glow(s, 4.5, 3.5, 8.0, RED, "0E")

# Sol: büyük slogan
txt(s, "Yalnız", 0.85, 0.95, 6.5, 1.0, 68, bold=True, color=WHITE)
txt(s, "gitme.", 0.85, 1.95, 6.5, 1.0, 68, bold=True, color=RED)

accent_line(s, 0.85, 3.1, 1.8, RED, 0.065)

txt(s,
    "Müzik zevki uyumluluğuna göre aynı\n"
    "etkinliğe gidecek biriyle eşleşiyorsunuz.",
    0.85, 3.28, 6.3, 0.9, 17, color=MUTED, italic=True)

# Özellik listesi
buddy_items = [
    ("🎯", "Uyumluluk skoru",      "0–100 tür eşleşme yüzdesi"),
    ("👆", "Kaydır eşleş",         "Geç → veya Birlikte Git ✓"),
    ("💫", "Eşleşme animasyonu",   "Ortak konser algılandığında konfeti"),
    ("🔗", "Profil keşfi",         "Eşleşen kişinin profiline git"),
]
for i, (ico, t_str, sub) in enumerate(buddy_items):
    lx, ty = 0.85, 4.38 + i*0.6
    card(s, lx, ty, 6.2, 0.52, BORDER, SURFACE2)
    txt(s, ico,   lx+0.15, ty+0.08, 0.5,  0.36, 16)
    txt(s, t_str, lx+0.8,  ty+0.08, 2.5,  0.36, 13, bold=True, color=WHITE)
    txt(s, sub,   lx+3.4,  ty+0.09, 2.7,  0.34, 12, color=MUTED)

# "Dünyada bir ilk" badge
pill_badge(s, 0.85, 6.9, 2.5, 0.38, RGBColor(0x00,0x2A,0x22), "★  DÜNYADA BİR İLK", 11, TEAL)

# Sağ: Buddy mockup
phone(s, "buddy", 8.35, 0.3, 3.8)

slide_number_dot(s, 6)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 7 — CONCERT PASSPORT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
glow(s, 4.0, 3.5, 9.0, PURPLE, "0D")

# Sol: Passport mockup
ph = phone(s, "passport", 0.6, 0.3, 3.6)

# Sağ metin
txt(s, "Her konser,", 5.5, 1.0, 7.5, 0.85, 50, bold=True, color=WHITE)
txt(s, "kalıcı.",     5.5, 1.85, 7.5, 0.85, 50, bold=True, color=PURPLE)

accent_line(s, 5.5, 2.88, 1.6, PURPLE, 0.065)

txt(s,
    "Kimin gittiğin, kim olduğunu anlatır.\n"
    "Doğrulanmış konser geçmişin, paylaşılabilir\n"
    "müzik kimliğine dönüşüyor.",
    5.5, 3.06, 7.3, 1.1, 17, color=MUTED, italic=True)

# İstatistik kartları
stats = [("🎟", "Konser", PURPLE), ("🎤", "Sanatçı", TEAL), ("📍", "Şehir", RED)]
for i, (ico, lab, col) in enumerate(stats):
    lx = 5.5 + i*2.45
    card(s, lx, 4.42, 2.2, 1.35, col, SURFACE2)
    txt(s, ico, lx+0.75, 4.58, 0.7, 0.55, 22, align=PP_ALIGN.CENTER)
    txt(s, lab, lx,      5.2,  2.2, 0.38, 12, bold=True, color=col, align=PP_ALIGN.CENTER)

# GPS doğrulama notu
card(s, 5.5, 6.05, 7.6, 0.92, TEAL, SURFACE2)
txt(s, "📍  GPS Doğrulama — Gerçekten orada olduğunu kanıtla.",
    5.7, 6.24, 7.2, 0.42, 14, bold=True, color=TEAL)

slide_number_dot(s, 7)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 8 — ETKİNLİK KEŞFİ + HARİTA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
glow(s, 7.0, 3.5, 10.0, TEAL, "09")

# Başlık
txt(s, "Biletinden haritasına,", 1.0, 0.3, 11.0, 0.7, 38, bold=True, color=WHITE)
txt(s, "tek platform.",          1.0, 1.0, 11.0, 0.7, 38, bold=True, color=TEAL)
accent_line(s, 1.0, 1.82, 1.8, TEAL, 0.055)

# İki mockup yan yana, ortada
pw = 3.3
gap = 0.65
tx = (W - 2*pw - gap) / 2

phone(s, "events", tx,         2.0, pw)
phone(s, "map_tr", tx+pw+gap,  2.0, pw)

# Altlarında etiket
txt(s, "99 etkinlik",    tx,          6.22, pw, 0.38, 14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txt(s, "986 etkinlik",   tx+pw+gap,   6.22, pw, 0.38, 14, bold=True, color=RED,  align=PP_ALIGN.CENTER)
txt(s, "Leman Sam · Paradise Lost · Sibel Can · Koray Avcı…",
    tx, 6.65, pw, 0.35, 10, color=MUTED, align=PP_ALIGN.CENTER)
txt(s, "Türkiye geneli · GPS filtreleme · 10–50 km yarıçap",
    tx+pw+gap, 6.65, pw, 0.35, 10, color=MUTED, align=PP_ALIGN.CENTER)

slide_number_dot(s, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 9 — TOPLULUK VE KİMLİK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
glow(s, 7.0, 3.5, 10.0, PURPLE, "0B")

txt(s, "Müzik kimliğin.", 0.9, 0.3, 11.0, 0.75, 44, bold=True, color=WHITE)
txt(s,
    "Türler · Sanatçılar · Topluluklar · Rozetler",
    0.9, 1.08, 11.0, 0.5, 20, color=MUTED)
accent_line(s, 0.9, 1.72, 2.0, PURPLE, 0.055)

# Üç mockup yan yana — fan düzeni
screens = [
    ("music",       "Müzik Profili",  PURPLE),
    ("communities", "Topluluklar",    RED),
    ("badges",      "Rozetler",       ORANGE),
]
pw3  = 2.7
gap3 = 0.65
tx3  = (W - 3*pw3 - 2*gap3) / 2
offsets = [0.15, 0.0, 0.15]  # orta mockup biraz yüksekte

for i, (key, label, col) in enumerate(screens):
    lx = tx3 + i*(pw3+gap3)
    ty = 1.95 + offsets[i]
    phone(s, key, lx, ty, pw3)
    # Kart etiketi
    pill_badge(s, lx+0.1, ty-0.45, pw3-0.2, 0.35,
               RGBColor(0x12,0x08,0x20), label, 11, col)

slide_number_dot(s, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 10 — REKABET POZİSYONU
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

txt(s, "Nerede duruyoruz?", 0.9, 0.28, 11.0, 0.62, 36, bold=True, color=WHITE)
txt(s, "Piyasada bu kesişim noktasını dolduran başka bir platform yok.",
    0.9, 0.94, 11.0, 0.45, 16, color=MUTED)
accent_line(s, 0.9, 1.52, 2.0, RED, 0.055)

# 2×2 Matris
MX, MY, MS = 1.2, 1.82, 5.0   # sol, üst, boyut
mid_l = MX + MS/2
mid_t = MY + MS/2

# Matris arka plan
rect(s, MX, MY, MS, MS, SURFACE2)

# Eksenler
rect(s, MX, mid_t-0.007, MS, 0.014, BORDER)
rect(s, mid_l-0.007, MY, 0.014, MS, BORDER)

# Eksen etiketleri
txt(s, "DÜŞÜK ETKİNLİK ODAĞI", MX, MY+MS+0.08, MS/2-0.1, 0.3, 9, color=MUTED)
txt(s, "YÜKSEK ETKİNLİK ODAĞI", MX+MS/2+0.1, MY+MS+0.08, MS/2, 0.3, 9, color=MUTED)
txt(s, "YÜKSEK SOSYAL", MX-1.1, MY+0.1, 1.0, MS/2, 11, color=MUTED)
txt(s, "DÜŞÜK SOSYAL",  MX-1.1, mid_t+0.1, 1.0, MS/2, 11, color=MUTED)

# Rakipler (koyu noktalar)
rivals = [
    ("Instagram",  mid_l-1.5, mid_t-0.6),
    ("WhatsApp",   mid_l-1.2, mid_t+0.4),
    ("Songkick",   mid_l+0.8, mid_t+0.5),
    ("Biletix",    mid_l+1.2, mid_t+0.8),
    ("Bandsintown",mid_l+0.5, mid_t+1.0),
]
for name, rx, ry in rivals:
    oval(s, rx-0.08, ry-0.08, 0.16, RGBColor(0x3A,0x3A,0x5A))
    txt(s, name, rx-0.5, ry+0.1, 1.0, 0.3, 9, color=MUTED)

# Concertly yıldızı — sol üst kadran (yüksek sosyal + yüksek etkinlik)
cx, cy = mid_l+1.1, mid_t-1.2
glow(s, cx, cy, 1.4, RED, "18")
oval(s, cx-0.18, cy-0.18, 0.36, RED)
txt(s, "★", cx-0.18, cy-0.2, 0.36, 0.36, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Concertly", cx-0.5, cy+0.22, 1.0, 0.3, 11, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Sağ taraf: karşılaştırma tablosu
tx2 = MX + MS + 0.5
tw  = W - tx2 - 0.4

card(s, tx2, 1.82, tw, 5.22, BORDER)
txt(s, "Özellik Karşılaştırması", tx2+0.25, 2.02, tw-0.4, 0.42, 14, bold=True, color=WHITE)

cols_header = ["Platform", "Buddy", "Passport", "GPS", "TR"]
cw = [2.0, 0.85, 1.0, 0.7, 0.6]
cx_starts = [tx2+0.2]
for i in range(1, len(cw)):
    cx_starts.append(cx_starts[-1]+cw[i-1]+0.05)

for j, h_str in enumerate(cols_header):
    col_c = RED if j == 0 else MUTED
    txt(s, h_str, cx_starts[j], 2.55, cw[j], 0.35, 11, bold=True, color=col_c)

rect(s, tx2+0.2, 2.94, tw-0.35, 0.012, BORDER)

rows = [
    ("Songkick",    "✗", "✗", "✗", "✗", MUTED),
    ("Bandsintown", "✗", "✗", "✗", "✗", MUTED),
    ("Instagram",   "✗", "✗", "✗", "✗", MUTED),
    ("Biletix",     "✗", "✗", "✗", "✓", MUTED),
    ("Concertly",   "✓", "✓", "✓", "✓", RED),
]
for ri, (platform, *checks, rc) in enumerate(rows):
    ty_r = 3.08 + ri*0.57
    is_us = (platform == "Concertly")
    if is_us:
        card(s, tx2+0.15, ty_r-0.05, tw-0.25, 0.52, RED, RGBColor(0x1F,0x08,0x12))
    txt(s, platform, cx_starts[0], ty_r, cw[0], 0.42, 12, bold=is_us, color=rc)
    for j, chk in enumerate(checks):
        chk_col = TEAL if chk == "✓" else RGBColor(0x44,0x44,0x60)
        txt(s, chk, cx_starts[j+1], ty_r, cw[j+1], 0.42, 13,
            bold=True, color=chk_col, align=PP_ALIGN.CENTER)

slide_number_dot(s, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 11 — YOL HARİTASI
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
glow(s, 7.0, 6.5, 8.0, TEAL, "09")

txt(s, "Nereden nereye?", 0.9, 0.28, 11.0, 0.65, 40, bold=True, color=WHITE)
txt(s, "Önce derinlik. Önce İstanbul. Sonra Türkiye.",
    0.9, 0.97, 11.0, 0.48, 18, color=MUTED)
accent_line(s, 0.9, 1.58, 2.0, TEAL, 0.055)

phases = [
    ("FAZ 1", "0–3 Ay",    "Stabilizasyon",  TEAL,
     ["Platform içi mesajlaşma", "Push bildirim (FCM/APNs)", "Güvenlik mekanizmaları",  "Onboarding iyileştirme"]),
    ("FAZ 2", "3–6 Ay",    "İçerik",         PURPLE,
     ["Organizatör hesabı",      "Self-servis etkinlik",     "Canlı Odalar (WebSocket)","Setlist paylaşımı"]),
    ("FAZ 3", "6–12 Ay",   "Ölçeklenme",     RED,
     ["Backend güçlendirme",     "ML öneri algoritması",     "Yeni şehirler",            "Web katmanı"]),
    ("FAZ 4", "12–24 Ay",  "Ekosistem",      ORANGE,
     ["Organizatör dashboard",   "Bilet entegrasyonu",       "Premium üyelik",           "Festival modu"]),
]

# Zaman çizgisi
line_y = 2.48
line_x0 = 1.1; line_x1 = W - 1.1
rect(s, line_x0, line_y-0.01, line_x1-line_x0, 0.025, BORDER)

pw4 = (line_x1 - line_x0) / len(phases)
for i, (faz, period, title, col, items) in enumerate(phases):
    lx = line_x0 + i*pw4
    cx = lx + pw4/2

    # Düğüm noktası
    glow(s, cx, line_y, 0.8, col, "18")
    oval(s, cx-0.16, line_y-0.16, 0.32, col)

    # Faz etiketi + süre
    txt(s, faz,    lx,      1.72, pw4, 0.38, 14, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(s, period, lx,      2.1,  pw4, 0.32, 11, color=MUTED, align=PP_ALIGN.CENTER)

    # Kart
    card_t = 2.82
    card(s, lx+0.12, card_t, pw4-0.24, 4.28, col, SURFACE2)
    txt(s, title, lx+0.22, card_t+0.18, pw4-0.42, 0.4, 14, bold=True, color=col)
    rect(s, lx+0.12, card_t+0.62, pw4-0.24, 0.012, BORDER)
    for j, item in enumerate(items):
        txt(s, f"▸  {item}", lx+0.22, card_t+0.78+j*0.74, pw4-0.38, 0.65, 11, color=WHITE)

# "Şu an buradayız" ok
txt(s, "◀ Şu an buradayız", line_x0+0.15, line_y+0.25, 2.5, 0.35, 10, color=TEAL)

slide_number_dot(s, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLAYT 12 — KAPANIŞ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# Büyük arka plan glow
glow(s, 4.0,  4.0, 12.0, RED,    "10")
glow(s, 10.0, 2.0, 8.0,  PURPLE, "0C")
glow(s, 7.0,  7.5, 7.0,  TEAL,   "0A")

# Sol dikey şerit
gradient_rect(s, 0, 0, 0.18, H, PURPLE, RED, angle=10800000)

# Logo
if os.path.exists(icon_path):
    s.shapes.add_picture(icon_path, Inches(5.92), Inches(0.6), Inches(1.5), Inches(1.5))
else:
    oval(s, 5.92, 0.6, 1.5, PURPLE)

# Karşılaştırma cümlesi
txt(s, "Spotify müziği nasıl dinlediğimizi değiştirdi.",
    1.5, 2.3, W-3.0, 0.62, 26, color=MUTED, align=PP_ALIGN.CENTER)

gradient_rect(s, 1.5, 3.05, W-3.0, 0.014, RED, PURPLE, angle=0)

txt(s, "Concertly, müziği nasıl yaşadığımızı değiştirmeye aday.",
    1.5, 3.22, W-3.0, 0.75, 30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Büyük kapanış sloganı
txt(s, "CONCERTLY",
    1.5, 4.18, W-3.0, 1.05, 62, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Gradient alt çizgi
gradient_rect(s, 3.5, 5.3, W-7.0, 0.08, RED, PURPLE, angle=0)

txt(s, "Müziği birlikte yaşa.",
    1.5, 5.52, W-3.0, 0.58, 22, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# Teknoloji pill'leri
tech_labels = [("React Native", PURPLE), ("Spring Boot", TEAL),
               ("Spotify API", RED), ("Ticketmaster", ORANGE)]
total_pw = len(tech_labels)*2.2 + (len(tech_labels)-1)*0.2
tx_pills = (W - total_pw) / 2
for i, (lab, col) in enumerate(tech_labels):
    pill_badge(s, tx_pills + i*2.4, 6.35, 2.1, 0.36,
               SURFACE2, lab, 10, col)

# Sorularınız?
txt(s, "Sorularınız?", 1.5, 6.9, W-3.0, 0.4,
    14, color=MUTED, align=PP_ALIGN.CENTER)

slide_number_dot(s, 12)


# ── Kaydet ────────────────────────────────────────────────────────────────────
output = r"C:\Users\EMRE\Desktop\Concertly_PitchDeck.pptx"
prs.save(output)
print(f"\nKaydedildi: {output}")
print(f"12 slayt, 16:9, {prs.slide_width.inches:.1f}x{prs.slide_height.inches:.1f} inc")
