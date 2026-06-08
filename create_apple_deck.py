"""
Concertly — Apple Keynote Tarzı Pitch Deck
Kural: Slayt başına maksimum 20 kelime. Tek mesaj. Büyük ürün.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palet ─────────────────────────────────────────────────────────────────────
BG      = RGBColor(0x08, 0x08, 0x12)   # neredeyse siyah
SURFACE = RGBColor(0x12, 0x12, 0x1E)
BORDER  = RGBColor(0x22, 0x22, 0x36)
FRAME   = RGBColor(0x18, 0x18, 0x28)   # telefon çerçevesi
SHADOW  = RGBColor(0x04, 0x04, 0x08)

RED     = RGBColor(0xE9, 0x45, 0x60)
TEAL    = RGBColor(0x00, 0xD4, 0xAA)
PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
ORANGE  = RGBColor(0xF5, 0xA6, 0x23)

WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x70, 0x70, 0x88)
DIM     = RGBColor(0x40, 0x40, 0x58)

# ── Kurulum ───────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H  = 13.33, 7.5

SCR_DIR = r"C:\Users\EMRE\Desktop\concertly-mobile\screenshots"
SCR = {
    "login":    SCR_DIR + r"\WhatsApp Image 2026-06-04 at 11.15.41.jpeg",
    "home":     SCR_DIR + r"\WhatsApp Image 2026-06-04 at 11.04.56.jpeg",
    "buddy":    SCR_DIR + r"\WhatsApp Image 2026-06-04 at 11.04.57 (2).jpeg",
    "passport": SCR_DIR + r"\WhatsApp Image 2026-06-04 at 11.04.56 (2).jpeg",
    "map":      SCR_DIR + r"\WhatsApp Image 2026-06-04 at 11.04.57 (3).jpeg",
    "events":   SCR_DIR + r"\WhatsApp Image 2026-06-04 at 11.04.58 (1).jpeg",
    "music":    SCR_DIR + r"\WhatsApp Image 2026-06-04 at 11.04.56 (3).jpeg",
    "badges":   SCR_DIR + r"\WhatsApp Image 2026-06-04 at 11.04.56 (1).jpeg",
    "comm":     SCR_DIR + r"\WhatsApp Image 2026-06-04 at 11.04.57.jpeg",
    "feed":     SCR_DIR + r"\WhatsApp Image 2026-06-04 at 11.04.58 (2).jpeg",
}

ICON = r"C:\Users\EMRE\Desktop\concertly-mobile\mobile\assets\icon.png"

# ── Temel araçlar ─────────────────────────────────────────────────────────────

def _rgb(c):
    h = str(c)
    return int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)

def apply_grad(sh, c1, c2, angle=0):
    sp   = sh._element
    _e   = sp.find(qn('p:spPr'))
    spPr = _e if _e is not None else etree.SubElement(sp, qn('p:spPr'))
    for tag in (qn('a:solidFill'), qn('a:gradFill')):
        for el in spPr.findall(tag):
            spPr.remove(el)
    gf  = etree.SubElement(spPr, qn('a:gradFill'))
    gsL = etree.SubElement(gf,  qn('a:gsLst'))
    for pos, rgb in [(0, _rgb(c1)), (100000, _rgb(c2))]:
        gs  = etree.SubElement(gsL, qn('a:gs'), pos=str(pos))
        etree.SubElement(gs, qn('a:srgbClr'), val=f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")
    etree.SubElement(gf, qn('a:lin'), ang=str(angle), scaled='0')

def set_alpha(sh, pct):
    """Shape'e saydamlık uygula (pct: 0-100)."""
    sp   = sh._element
    _e   = sp.find(qn('p:spPr'))
    spPr = _e if _e is not None else etree.SubElement(sp, qn('p:spPr'))
    sf   = spPr.find(qn('a:solidFill'))
    if sf is not None:
        c = sf.find(qn('a:srgbClr'))
        if c is not None:
            a = etree.SubElement(c, qn('a:alpha'))
            a.set('val', str(int(pct * 1000)))

def bg(s, color=BG):
    sh = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def box(s, l, t, w, h, color=SURFACE, alpha=None):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    if alpha is not None:
        set_alpha(sh, alpha)
    return sh

def pill(s, l, t, w, h, color, alpha=None):
    sh = s.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    if alpha is not None:
        set_alpha(sh, alpha)
    return sh

def oval(s, l, t, w, h, color, alpha=None):
    sh = s.shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    if alpha is not None:
        set_alpha(sh, alpha)
    return sh

def grad_box(s, l, t, w, h, c1, c2, angle=0):
    sh = box(s, l, t, w, h, c1)
    apply_grad(sh, c1, c2, angle)
    return sh

def label(s, text, l, t, w, h, size, bold=False,
          color=WHITE, align=PP_ALIGN.LEFT, italic=False, font="Inter"):
    b  = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = b.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.name   = font
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return b

def glow(s, cx, cy, size, color, alpha=8):
    """Büyük, saydam radyal ışıma."""
    oval(s, cx-size/2, cy-size/2, size, size, color, alpha)

def phone(s, key, l, t, w, reflect=True):
    """
    iPhone mockup — büyük, temiz, Apple tarzı.
    w: genişlik (inç), yükseklik otomatik hesaplanır (2048/945 oranı).
    """
    ratio = 2048 / 945
    h     = w * ratio

    # Altta yansıma / glow efekti
    if reflect:
        glow(s, l+w/2, t+h+0.3, w*1.6, color=WHITE, alpha=4)

    # Arka gölge
    shadow = s.shapes.add_shape(
        5, Inches(l+0.12), Inches(t+0.14), Inches(w), Inches(h))
    shadow.fill.solid(); shadow.fill.fore_color.rgb = SHADOW
    shadow.line.fill.background()

    # Çerçeve (dışarıda hafif çıkıntı)
    frame = s.shapes.add_shape(
        5, Inches(l-0.09), Inches(t-0.09), Inches(w+0.18), Inches(h+0.18))
    frame.fill.solid(); frame.fill.fore_color.rgb = FRAME
    frame.line.color.rgb = BORDER; frame.line.width = Inches(0.022)

    # Görüntü
    path = SCR.get(key, "")
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        fb = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        fb.fill.solid(); fb.fill.fore_color.rgb = SURFACE; fb.line.fill.background()
        label(s, key, l, t+h/2-0.2, w, 0.4, 11, color=MUTED, align=PP_ALIGN.CENTER)
    return h

def slide_no(s, n, total=12):
    label(s, f"{n}  /  {total}", W-1.4, H-0.42, 1.2, 0.35,
          10, color=DIM, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  S1 — KAPAK
#  Mesaj: "CONCERTLY · Müziği birlikte yaşa." (4 kelime)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

# Arka plan glowlar
glow(s,  2.5, 5.5, 9.0, RED,    5)
glow(s, 10.5, 2.0, 6.0, PURPLE, 5)

# Logo
if os.path.exists(ICON):
    s.shapes.add_picture(ICON, Inches(1.1), Inches(1.55), Inches(1.1), Inches(1.1))
else:
    oval(s, 1.1, 1.55, 1.1, 1.1, PURPLE)

# Başlık
label(s, "CONCERTLY", 1.1, 2.8, 7.5, 1.5, 90, bold=True, color=WHITE)

# Gradient çizgi
grad_box(s, 1.1, 4.38, 5.5, 0.07, RED, PURPLE, angle=0)

# Tagline
label(s, "Müziği birlikte yaşa.", 1.1, 4.58, 7.0, 0.65, 26, italic=True, color=MUTED)

# Sağda büyük telefon
phone(s, "login", 8.6, 0.25, 3.85)

slide_no(s, 1)

# ══════════════════════════════════════════════════════════════════════════════
#  S2 — O AN
#  Mesaj: "Konsere müzik için değil, o an için gidiyorsunuz." (9 kelime)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

# Büyük kırmızı glow ortada
glow(s, W/2, H/2, 14.0, RED, 6)
glow(s, W/2, H/2,  8.0, RED, 8)

# Üst küçük metin
label(s, "Neden konsere gidiyoruz?",
      0, 1.45, W, 0.55, 22, color=MUTED, align=PP_ALIGN.CENTER)

# Büyük mesaj — 2 satır
label(s, "Müzik için değil.",
      0, 2.2, W, 1.1, 72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
label(s, "O an için.",
      0, 3.3, W, 1.1, 72, bold=True, color=RED,   align=PP_ALIGN.CENTER)

# Alt ince not
label(s, "Konserdeki o titreşim — bir daha olmayacak.",
      0, 4.85, W, 0.5, 18, italic=True, color=DIM, align=PP_ALIGN.CENTER)

slide_no(s, 2)

# ══════════════════════════════════════════════════════════════════════════════
#  S3 — PROBLEM
#  Mesaj: "5 platform. 1 gece. 0 iz." (6 kelime)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

glow(s, W/2, 3.2, 10.0, PURPLE, 5)

label(s, "Bugün bir konser planlamak için:",
      0, 0.5, W, 0.55, 20, color=MUTED, align=PP_ALIGN.CENTER)

# 5 ikon akışı
icons  = ["🎟", "📸", "💬", "📅", "⏱"]
labels = ["Biletix", "Instagram", "WhatsApp", "Takvim", "24 saat"]
colors = [PURPLE, PURPLE, PURPLE, PURPLE, RED]
iw     = 1.8
gap    = 0.55
total  = len(icons)*iw + (len(icons)-1)*gap
x0     = (W - total) / 2
iy     = 1.3

for i, (ic, lb, col) in enumerate(zip(icons, labels, colors)):
    lx = x0 + i*(iw+gap)
    is_last = i == len(icons)-1

    # Kart
    sh = s.shapes.add_shape(5, Inches(lx), Inches(iy), Inches(iw), Inches(1.75))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0x1A,0x08,0x14) if is_last else SURFACE
    sh.line.color.rgb = RED if is_last else BORDER
    sh.line.width = Inches(0.018 if is_last else 0.012)

    label(s, ic, lx, iy+0.18, iw, 0.75, 30, align=PP_ALIGN.CENTER)
    label(s, lb, lx, iy+0.98, iw, 0.38, 12,
          bold=True, color=RED if is_last else WHITE, align=PP_ALIGN.CENTER)

    # Ok
    if not is_last:
        label(s, "→", lx+iw+0.12, iy+0.55, 0.38, 0.55,
              22, color=DIM, align=PP_ALIGN.CENTER)

# Ana mesaj — büyük
label(s, "5 platform.",
      0, 3.52, W*0.33, 0.95, 54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
label(s, "1 gece.",
      W*0.33, 3.52, W*0.33, 0.95, 54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
label(s, "0 iz.",
      W*0.66, 3.52, W*0.33, 0.95, 54, bold=True, color=RED,   align=PP_ALIGN.CENTER)

# Ince divider
box(s, 1.0, 3.45, W-2.0, 0.012, BORDER)

# Problem badges
probs = [
    ("Etkinlik keşfi dağınık", PURPLE),
    ("Yalnız gitme korkusu",   RED),
    ("Konser geçmişi kayıt yok", TEAL),
]
pw = 3.5; pg = 0.3
px = (W - 3*pw - 2*pg) / 2
for i, (txt_str, col) in enumerate(probs):
    lx = px + i*(pw+pg)
    sh = s.shapes.add_shape(5, Inches(lx), Inches(4.78), Inches(pw), Inches(0.62))
    sh.fill.solid(); sh.fill.fore_color.rgb = SURFACE; sh.line.fill.background()
    label(s, txt_str, lx, 4.88, pw, 0.42, 14, bold=True, color=col, align=PP_ALIGN.CENTER)

label(s, "Bu deneyim dağınık — bir arada çözülmedi.",
      0, 5.7, W, 0.45, 15, italic=True, color=DIM, align=PP_ALIGN.CENTER)

slide_no(s, 3)

# ══════════════════════════════════════════════════════════════════════════════
#  S4 — VİZYON
#  Mesaj: "Canlı müziğin dijital katmanı." (4 kelime)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

glow(s, W/2, H/2, 15.0, PURPLE, 5)
glow(s, W/2, H/2,  7.0, RED,    6)

label(s, "Canlı müziğin",
      0, 1.6, W, 1.25, 82, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Gradient kutu altında metin
grd = grad_box(s, 1.8, 2.92, W-3.6, 1.12, RED, PURPLE, angle=0)
label(s, "dijital katmanı.",
      1.8, 2.92, W-3.6, 1.12, 82, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

label(s, "Önce · Sırasında · Sonra",
      0, 4.4, W, 0.55, 20, color=DIM, align=PP_ALIGN.CENTER)

slide_no(s, 4)

# ══════════════════════════════════════════════════════════════════════════════
#  S5 — KEŞFET
#  Mesaj: "Her etkinlik. Tek platform." (4 kelime)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

glow(s, 3.8, 3.5, 9.0, PURPLE, 6)

# Sol metin bloğu
label(s, "Her etkinlik.",   0.9, 1.4, 6.5, 1.05, 62, bold=True, color=WHITE)
label(s, "Tek platform.",   0.9, 2.45, 6.5, 1.05, 62, bold=True, color=PURPLE)

grad_box(s, 0.9, 3.65, 2.0, 0.07, RED, PURPLE, angle=0)

label(s, "Şehir · Tür · Sanatçı · Tarih", 0.9, 3.88, 6.0, 0.5, 17, color=MUTED)

# Stat badge
sh = s.shapes.add_shape(5, Inches(0.9), Inches(4.7), Inches(5.2), Inches(0.72))
sh.fill.solid(); sh.fill.fore_color.rgb = SURFACE; sh.line.fill.background()
label(s, "99 etkinlik  ·  11 şehir  ·  Türkiye geneli",
      0.9, 4.88, 5.2, 0.38, 15, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

# Sağda büyük telefon
phone(s, "home", 8.4, 0.2, 3.95)

slide_no(s, 5)

# ══════════════════════════════════════════════════════════════════════════════
#  S6 — CONCERT BUDDY
#  Mesaj: "Yalnız gitme." (2 kelime)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

glow(s,  4.0, 3.5, 10.0, RED, 6)
glow(s, 10.5, 3.0,  5.0, RED, 5)

# Sol telefon — büyük
phone(s, "buddy", 0.45, 0.2, 4.1)

# Sağ metin
label(s, "Yalnız",  6.3, 1.2, 6.8, 1.25, 88, bold=True, color=WHITE)
label(s, "gitme.",  6.3, 2.45, 6.8, 1.25, 88, bold=True, color=RED)

grad_box(s, 6.3, 3.85, 2.2, 0.07, RED, PURPLE, angle=0)

label(s, "Müzik zevki uyumluluğuna göre\naynı etkinliğe gidecek biriyle eşleş.",
      6.3, 4.05, 6.8, 0.95, 18, color=MUTED, italic=True)

# "Dünyada ilk" badge
sh = s.shapes.add_shape(5, Inches(6.3), Inches(5.55), Inches(3.2), Inches(0.55))
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x00,0x22,0x1A); sh.line.fill.background()
label(s, "★  DÜNYADA BİR İLK", 6.3, 5.64, 3.2, 0.35, 12, bold=True,
      color=TEAL, align=PP_ALIGN.CENTER)

slide_no(s, 6)

# ══════════════════════════════════════════════════════════════════════════════
#  S7 — CONCERT PASSPORT
#  Mesaj: "Her konser, kalıcı." (3 kelime)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

glow(s,  9.5, 3.5, 10.0, PURPLE, 6)
glow(s, 11.5, 5.0,  5.0, TEAL,   5)

# Sol metin
label(s, "Her konser,", 0.9, 1.2, 6.5, 1.1, 72, bold=True, color=WHITE)
label(s, "kalıcı.",     0.9, 2.3, 6.5, 1.1, 72, bold=True, color=PURPLE)

grad_box(s, 0.9, 3.55, 2.0, 0.07, PURPLE, TEAL, angle=0)

label(s, "Doğrulanmış · Paylaşılabilir · Sonsuz",
      0.9, 3.75, 6.5, 0.5, 18, color=MUTED)

# 3 istatistik
stats = [("🎟", "Konser", PURPLE), ("🎤", "Sanatçı", TEAL), ("📍", "Şehir", RED)]
for i, (ico, lab, col) in enumerate(stats):
    lx = 0.9 + i*2.15
    sh = s.shapes.add_shape(5, Inches(lx), Inches(4.55), Inches(1.9), Inches(1.3))
    sh.fill.solid(); sh.fill.fore_color.rgb = SURFACE; sh.line.fill.background()
    label(s, ico, lx, 4.68, 1.9, 0.55, 22, align=PP_ALIGN.CENTER)
    label(s, lab, lx, 5.26, 1.9, 0.36, 12, bold=True, color=col, align=PP_ALIGN.CENTER)

# GPS notu
label(s, "GPS doğrulama ile gerçek katılım kanıtla.",
      0.9, 6.18, 6.5, 0.45, 14, color=DIM, italic=True)

# Sağ telefon
phone(s, "passport", 8.35, 0.2, 4.0)

slide_no(s, 7)

# ══════════════════════════════════════════════════════════════════════════════
#  S8 — HARİTA
#  Mesaj: "986 etkinlik. Türkiye geneli." (4 kelime)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

glow(s, W/2, 3.2, 14.0, TEAL, 5)

label(s, "986 etkinlik.",
      0, 0.45, W, 1.0, 72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
label(s, "Türkiye geneli.",
      0, 1.45, W, 1.0, 72, bold=True, color=TEAL,  align=PP_ALIGN.CENTER)

# Büyük harita telefonu — merkeze yakın
ph = phone(s, "map", 3.85, 1.55, 5.6)

# Alt bilgi
label(s, "GPS filtreleme  ·  10 – 25 – 50 km yarıçap  ·  Anlık konum",
      0, H-0.55, W, 0.4, 13, color=DIM, align=PP_ALIGN.CENTER)

slide_no(s, 8)

# ══════════════════════════════════════════════════════════════════════════════
#  S9 — KİMLİK
#  Mesaj: "Müzik kimliğin." (2 kelime)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

glow(s, W/2, 3.0, 12.0, PURPLE, 5)

label(s, "Müzik kimliğin.",
      0, 0.38, W, 0.95, 58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

label(s, "Türler  ·  Sanatçılar  ·  Topluluklar  ·  Rozetler",
      0, 1.38, W, 0.5, 18, color=MUTED, align=PP_ALIGN.CENTER)

# 3 telefon — hafif eğimli fan düzeni
screens = [
    ("music", -3.5, PURPLE),
    ("comm",   0.0, RED),
    ("badges", 3.5, ORANGE),
]
pw3 = 3.05
ph3 = pw3 * (2048/945)
centers = [3.22, 6.64, 10.06]
tops    = [2.05, 1.72, 2.05]

for i, ((key, rot, col), cx, ty) in enumerate(zip(screens, centers, tops)):
    lx = cx - pw3/2
    phone(s, key, lx, ty, pw3)
    # Etiket pill
    sh = s.shapes.add_shape(5, Inches(lx+0.15), Inches(ty-0.5),
                             Inches(pw3-0.3), Inches(0.38))
    sh.fill.solid()
    sh.fill.fore_color.rgb = SURFACE
    sh.line.fill.background()
    names = ["Müzik Profili", "Topluluklar", "Rozetler"]
    label(s, names[i], lx+0.15, ty-0.44, pw3-0.3, 0.34, 11,
          bold=True, color=col, align=PP_ALIGN.CENTER)

slide_no(s, 9)

# ══════════════════════════════════════════════════════════════════════════════
#  S10 — FARK
#  Mesaj: "Dünyada benzeri yok." (3 kelime)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

glow(s, W/2, H/2, 12.0, RED, 5)

label(s, "Dünyada benzeri yok.",
      0, 0.32, W, 0.85, 52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
label(s, "Concert Buddy · Concert Passport · GPS Doğrulama",
      0, 1.2, W, 0.5, 18, color=MUTED, align=PP_ALIGN.CENTER)

# Tablo — minimal, temiz
PLATFORMS = [
    ("Songkick",    False, False, False, False),
    ("Bandsintown", False, False, False, False),
    ("Instagram",   False, False, False, False),
    ("Biletix",     False, False, False, True),
    ("Concertly",   True,  True,  True,  True),
]
COLS     = ["", "Concert Buddy", "Passport", "GPS Doğr.", "TR Odağı"]
COL_W    = [2.6, 2.3, 2.0, 1.9, 1.9]
TABLE_W  = sum(COL_W) + 0.2
TX       = (W - TABLE_W) / 2
TY       = 2.0
ROW_H    = 0.7

# Header
for j, (col_txt, cw) in enumerate(zip(COLS, COL_W)):
    cx = TX + sum(COL_W[:j]) + j*0.04
    label(s, col_txt, cx, TY, cw, 0.45, 12,
          bold=True, color=MUTED, align=PP_ALIGN.CENTER if j else PP_ALIGN.LEFT)

box(s, TX, TY+0.48, TABLE_W, 0.012, BORDER)

for ri, (plat, *checks) in enumerate(PLATFORMS):
    ry = TY + 0.55 + ri*ROW_H
    is_us = plat == "Concertly"

    if is_us:
        grad_box(s, TX-0.1, ry+0.04, TABLE_W+0.2, ROW_H-0.06,
                 RGBColor(0x1A,0x06,0x10), RGBColor(0x12,0x04,0x1E), angle=0)

    cx = TX
    plat_col = RED if is_us else MUTED
    label(s, plat, cx, ry+0.14, COL_W[0], 0.42,
          15 if is_us else 14, bold=is_us, color=plat_col)

    for j, ok in enumerate(checks):
        cx = TX + sum(COL_W[:j+1]) + (j+1)*0.04
        sym  = "✓" if ok else "–"
        col  = TEAL if ok else DIM
        size = 20 if ok else 16
        label(s, sym, cx, ry+0.12, COL_W[j+1], 0.46,
              size, bold=ok, color=col, align=PP_ALIGN.CENTER)

    box(s, TX, ry+ROW_H-0.06, TABLE_W, 0.008, RGBColor(0x18,0x18,0x28))

slide_no(s, 10)

# ══════════════════════════════════════════════════════════════════════════════
#  S11 — YOL HARİTASI
#  Mesaj: "Önce derinlik. Sonra genişlik." (4 kelime)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

glow(s, W/2, H/2, 12.0, PURPLE, 5)

label(s, "Önce derinlik.",   0.9, 0.35, W-1.8, 0.85, 52, bold=True, color=WHITE)
label(s, "Sonra genişlik.",  0.9, 1.2,  W-1.8, 0.85, 52, bold=True, color=PURPLE)
label(s, "Önce İstanbul. Sonra Türkiye.",
      0.9, 2.1, W-1.8, 0.48, 18, color=MUTED)

PHASES = [
    ("FAZ 1",  "0–3 Ay",   "Stabilizasyon", TEAL,
     ["Mesajlaşma", "Push bildirim", "Güvenlik", "Onboarding"]),
    ("FAZ 2",  "3–6 Ay",   "İçerik",        PURPLE,
     ["Organizatör", "Canlı Odalar", "Setlist", "Yerel etkinlik"]),
    ("FAZ 3",  "6–12 Ay",  "Ölçeklenme",    RED,
     ["Backend", "ML öneri", "Yeni şehirler", "Web"]),
    ("FAZ 4",  "12–24 Ay", "Ekosistem",     ORANGE,
     ["Bilet entegr.", "Premium", "Dashboard", "Festival modu"]),
]

LINE_Y = 2.92
X0, X1 = 1.0, W-1.0
box(s, X0, LINE_Y-0.01, X1-X0, 0.022, BORDER)

PW = (X1-X0) / len(PHASES)
for i, (faz, period, title, col, items) in enumerate(PHASES):
    cx = X0 + i*PW + PW/2

    # Düğüm — glow + daire
    glow(s, cx, LINE_Y, 1.2, col, 12)
    oval(s, cx-0.2, LINE_Y-0.2, 0.4, 0.4, col)

    # Üstte faz + süre
    label(s, faz,    X0+i*PW, LINE_Y-0.82, PW, 0.35, 12, bold=True, color=col, align=PP_ALIGN.CENTER)
    label(s, period, X0+i*PW, LINE_Y-0.44, PW, 0.3,  11, color=DIM,  align=PP_ALIGN.CENTER)

    # Altta kart
    CY = LINE_Y + 0.45
    CW = PW - 0.35
    CL = X0 + i*PW + 0.18
    sh = s.shapes.add_shape(
        1, Inches(CL), Inches(CY), Inches(CW), Inches(3.72))
    sh.fill.solid(); sh.fill.fore_color.rgb = SURFACE; sh.line.fill.background()

    # Renkli üst çizgi
    box(s, CL, CY, CW, 0.06, col)

    label(s, title, CL, CY+0.14, CW, 0.42, 15, bold=True, color=col, align=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        label(s, item, CL+0.18, CY+0.68+j*0.7, CW-0.25, 0.62, 13, color=WHITE)

# Şu an oku
label(s, "Şu an buradayız  ▶",
      X0-0.05, LINE_Y+0.26, 2.2, 0.38, 11, color=TEAL)

slide_no(s, 11)

# ══════════════════════════════════════════════════════════════════════════════
#  S12 — KAPANIŞ
#  Mesaj: "CONCERTLY · Müziği birlikte yaşa." (4 kelime)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

glow(s, W/2, 3.2, 15.0, RED,    6)
glow(s, W/2, 4.5, 10.0, PURPLE, 7)
glow(s, W/2, 6.5,  8.0, TEAL,   5)

# Logo
if os.path.exists(ICON):
    s.shapes.add_picture(ICON, Inches(W/2-0.7), Inches(0.55), Inches(1.4), Inches(1.4))
else:
    oval(s, W/2-0.7, 0.55, 1.4, 1.4, PURPLE)

# Spotify karşılaştırma — küçük, üstte
label(s, "Spotify müziği nasıl dinlediğimizi değiştirdi.",
      0, 2.12, W, 0.52, 20, color=DIM, align=PP_ALIGN.CENTER)

grad_box(s, 2.5, 2.72, W-5.0, 0.012, RED, PURPLE, angle=0)

label(s, "Concertly, müziği nasıl yaşadığımızı değiştirmeye aday.",
      0, 2.85, W, 0.62, 22, color=MUTED, align=PP_ALIGN.CENTER)

# Büyük final
label(s, "CONCERTLY",
      0, 3.6, W, 1.5, 100, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

grad_box(s, 3.2, 5.18, W-6.4, 0.1, RED, PURPLE, angle=0)

label(s, "Müziği birlikte yaşa.",
      0, 5.42, W, 0.65, 26, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# Teknoloji pill'leri
TECHS = [("React Native", PURPLE), ("Spring Boot", TEAL), ("Spotify API", RED), ("Ticketmaster", ORANGE)]
PW4 = 2.1; PG = 0.22
ptotal = len(TECHS)*PW4 + (len(TECHS)-1)*PG
px0 = (W - ptotal) / 2
for i, (t, c) in enumerate(TECHS):
    lx = px0 + i*(PW4+PG)
    sh = s.shapes.add_shape(5, Inches(lx), Inches(6.38), Inches(PW4), Inches(0.42))
    sh.fill.solid(); sh.fill.fore_color.rgb = SURFACE; sh.line.fill.background()
    label(s, t, lx, 6.46, PW4, 0.3, 11, bold=True, color=c, align=PP_ALIGN.CENTER)

label(s, "Sorularınız?",
      0, H-0.44, W, 0.36, 13, color=DIM, align=PP_ALIGN.CENTER)

slide_no(s, 12)

# ── Kaydet ────────────────────────────────────────────────────────────────────
OUT = r"C:\Users\EMRE\Desktop\Concertly_Apple_Deck.pptx"
prs.save(OUT)
print(f"Kaydedildi: {OUT}")
print(f"12 slayt  |  16:9  |  {os.path.getsize(OUT)//1024} KB")
