from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette — Koyu: Siyah / Mor / Lacivert / Turuncu ──────────────────────
BG      = RGBColor(0x0A, 0x08, 0x1A)   # siyah-mor arka plan
BG2     = RGBColor(0x0E, 0x0C, 0x22)   # biraz daha açık siyah-mor
CARD    = RGBColor(0x14, 0x12, 0x2E)   # koyu lacivert-mor kart
CARD2   = RGBColor(0x1C, 0x18, 0x3C)   # lacivert kart
CARD_DK = RGBColor(0x08, 0x06, 0x14)   # neredeyse siyah

PURPLE  = RGBColor(0x7C, 0x3A, 0xED)   # canlı mor
BLUE    = RGBColor(0x1D, 0x4E, 0xD8)   # koyu kobalt mavi / lacivert
ORANGE  = RGBColor(0xF5, 0xA6, 0x23)   # turuncu
RED     = RGBColor(0xE9, 0x45, 0x60)   # canlı kırmızı-koral

WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0x8A, 0x8A, 0xAA)   # mor tonlu gri
DGRAY   = RGBColor(0x22, 0x20, 0x40)   # koyu lacivert-gri

TOTAL = 17

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = 13.33, 7.5

# ── XML helper: sadece gradient için kullanılır ────────────────────────────
def apply_grad(shape, stops, ang=5400000):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(sp, qn('p:spPr'))
    for old in spPr.findall(qn('a:solidFill')) + spPr.findall(qn('a:gradFill')):
        spPr.remove(old)
    gf = etree.SubElement(spPr, qn('a:gradFill'))
    gsLst = etree.SubElement(gf, qn('a:gsLst'))
    for pos, rgb in stops:
        gs = etree.SubElement(gsLst, qn('a:gs'), pos=str(int(pos * 100000)))
        etree.SubElement(gs, qn('a:srgbClr'), val=str(rgb))
    etree.SubElement(gf, qn('a:lin'), ang=str(ang), scaled='0')

# ── Temel çizim fonksiyonları ──────────────────────────────────────────────

def slide_bg(s):
    """Solid siyah-mor arka plan — gradient yok, güvenli"""
    sh = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = BG; sh.line.fill.background()
    return sh

def rect(s, l, t, w, h, color=CARD):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    return sh

def card(s, l, t, w, h, border_color=None):
    """Gölgeli kart"""
    sd = s.shapes.add_shape(1, Inches(l+0.07), Inches(t+0.08), Inches(w), Inches(h))
    sd.fill.solid(); sd.fill.fore_color.rgb = CARD_DK; sd.line.fill.background()
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = CARD
    if border_color:
        sh.line.color.rgb = border_color; sh.line.width = Inches(0.018)
    else:
        sh.line.fill.background()
    return sh

def card2(s, l, t, w, h, border_color=None):
    """CARD2 rengiyle kart"""
    sd = s.shapes.add_shape(1, Inches(l+0.07), Inches(t+0.08), Inches(w), Inches(h))
    sd.fill.solid(); sd.fill.fore_color.rgb = CARD_DK; sd.line.fill.background()
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = CARD2
    if border_color:
        sh.line.color.rgb = border_color; sh.line.width = Inches(0.018)
    else:
        sh.line.fill.background()
    return sh

def circle(s, l, t, d, color=PURPLE):
    sh = s.shapes.add_shape(9, Inches(l), Inches(t), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    return sh

def txt(s, text, l, t, w, h, size, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    b.word_wrap = True
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return b

def aline(s, x, y, w=1.4, color=PURPLE, thick=0.07):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(thick))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

def pill(s, l, t, w, h, color, label, lsize=11, lcolor=WHITE):
    sh = s.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    txt(s, label, l, t+0.01, w, h-0.02, lsize, bold=True, color=lcolor, align=PP_ALIGN.CENTER)

def pbar(s, n, color=PURPLE):
    """Üstte progress bar + altta dot göstergesi"""
    tr = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.05))
    tr.fill.solid(); tr.fill.fore_color.rgb = DGRAY; tr.line.fill.background()
    bar = s.shapes.add_shape(1, 0, 0, Inches(W * n / TOTAL), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    dot_total_w = TOTAL * 0.16
    start = W / 2 - dot_total_w / 2
    for i in range(TOTAL):
        cx = start + i * 0.16 + 0.035
        d = s.shapes.add_shape(9, Inches(cx), Inches(H-0.2), Inches(0.09), Inches(0.09))
        d.fill.solid()
        d.fill.fore_color.rgb = color if i < n else DGRAY
        d.line.fill.background()

# ── Ekran görüntüsü yolları ────────────────────────────────────────────────
_SCR = r"C:\Users\EMRE\Desktop\concertly-mobile\screenshots"
SCREEN_HOME    = _SCR + r"\WhatsApp Image 2026-06-04 at 11.04.56.jpeg"
SCREEN_BADGES  = _SCR + r"\WhatsApp Image 2026-06-04 at 11.04.56 (1).jpeg"
SCREEN_PASSPORT= _SCR + r"\WhatsApp Image 2026-06-04 at 11.04.56 (2).jpeg"
SCREEN_MUSIC   = _SCR + r"\WhatsApp Image 2026-06-04 at 11.04.56 (3).jpeg"
SCREEN_COMMUNITIES = _SCR + r"\WhatsApp Image 2026-06-04 at 11.04.57.jpeg"
SCREEN_PROFILE = _SCR + r"\WhatsApp Image 2026-06-04 at 11.04.57 (1).jpeg"
SCREEN_BUDDY   = _SCR + r"\WhatsApp Image 2026-06-04 at 11.04.57 (2).jpeg"
SCREEN_MAP_TR  = _SCR + r"\WhatsApp Image 2026-06-04 at 11.04.57 (3).jpeg"
SCREEN_MAP_ANK = _SCR + r"\WhatsApp Image 2026-06-04 at 11.04.58.jpeg"
SCREEN_EVENTS  = _SCR + r"\WhatsApp Image 2026-06-04 at 11.04.58 (1).jpeg"
SCREEN_FEED    = _SCR + r"\WhatsApp Image 2026-06-04 at 11.04.58 (2).jpeg"
SCREEN_REGISTER= _SCR + r"\WhatsApp Image 2026-06-04 at 11.15.40.jpeg"
SCREEN_LOGIN   = _SCR + r"\WhatsApp Image 2026-06-04 at 11.15.41.jpeg"
SCREEN_ADMIN   = _SCR + r"\WhatsApp Image 2026-06-04 at 11.15.41 (1).jpeg"

ALL_SCREENS = [
    (SCREEN_HOME,         "Ana Sayfa"),
    (SCREEN_EVENTS,       "Etkinlikler"),
    (SCREEN_FEED,         "Feed"),
    (SCREEN_COMMUNITIES,  "Topluluklar"),
    (SCREEN_BUDDY,        "Konser Arkadaşı"),
    (SCREEN_PASSPORT,     "Konser Pasaportu"),
    (SCREEN_MAP_TR,       "Harita (TR)"),
    (SCREEN_MAP_ANK,      "Harita (Ankara)"),
    (SCREEN_PROFILE,      "Profil"),
    (SCREEN_BADGES,       "Rozetler"),
    (SCREEN_MUSIC,        "Müzik Profili"),
]

def phone_screenshot(s, img_path, l, t, w):
    """Gerçek ekran görüntüsü — iPhone oranı 945:2048"""
    h = w * 2.167
    sd = s.shapes.add_shape(1, Inches(l+0.07), Inches(t+0.09),
                             Inches(w+0.04), Inches(h+0.04))
    sd.fill.solid(); sd.fill.fore_color.rgb = CARD_DK; sd.line.fill.background()
    frame = s.shapes.add_shape(5, Inches(l-0.06), Inches(t-0.06),
                                Inches(w+0.12), Inches(h+0.12))
    frame.fill.solid(); frame.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x30)
    frame.line.color.rgb = DGRAY; frame.line.width = Inches(0.015)
    try:
        s.shapes.add_picture(img_path, Inches(l), Inches(t), Inches(w), Inches(h))
    except Exception:
        ph = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        ph.fill.solid(); ph.fill.fore_color.rgb = CARD2; ph.line.fill.background()
    return h

def slide_header(s, title, subtitle, accent=PURPLE):
    """Başlık alanı: solid koyu şerit + çizgi"""
    # Başlık arkasına koyu solid şerit
    strip = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.9))
    strip.fill.solid(); strip.fill.fore_color.rgb = CARD2; strip.line.fill.background()
    # Sol renk çubuğu
    aline(s, 0.7, 0.35, 1.6, accent, 0.08)
    txt(s, title,    0.7, 0.46, 11.5, 0.88, 42, bold=True, color=WHITE)
    txt(s, subtitle, 0.7, 1.37, 11,   0.42, 15, color=LGRAY)
    # Ayırıcı çizgi
    div = s.shapes.add_shape(1, Inches(0.7), Inches(1.83), Inches(W-1.4), Inches(0.015))
    div.fill.solid(); div.fill.fore_color.rgb = accent; div.line.fill.background()

def phone_mockup(s, l, t, w, content_lines=None, header=None,
                 header_color=PURPLE, bg_color=CARD2):
    h = w * 2.1
    body = s.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    body.fill.solid(); body.fill.fore_color.rgb = bg_color
    body.line.color.rgb = DGRAY; body.line.width = Inches(0.025)
    notch = s.shapes.add_shape(9, Inches(l+w*0.35), Inches(t+0.07),
                               Inches(w*0.3), Inches(0.12))
    notch.fill.solid(); notch.fill.fore_color.rgb = DGRAY; notch.line.fill.background()
    circle(s, l+0.08, t+0.06, 0.08, DGRAY)
    if header:
        hbar = s.shapes.add_shape(1, Inches(l+0.04), Inches(t+0.27),
                                  Inches(w-0.08), Inches(0.32))
        hbar.fill.solid(); hbar.fill.fore_color.rgb = header_color; hbar.line.fill.background()
        txt(s, header, l+0.08, t+0.28, w-0.16, 0.28, 8, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
    if content_lines:
        cy = t + 0.65
        for (line_txt, lsize, lcolor) in content_lines:
            if line_txt == "---":
                sep = s.shapes.add_shape(1, Inches(l+0.06), Inches(cy),
                                         Inches(w-0.12), Inches(0.02))
                sep.fill.solid(); sep.fill.fore_color.rgb = DGRAY; sep.line.fill.background()
                cy += 0.1
            else:
                txt(s, line_txt, l+0.1, cy, w-0.2, 0.28, lsize, color=lcolor)
                cy += 0.28
    ind = s.shapes.add_shape(1, Inches(l+w*0.3), Inches(t+h-0.12),
                             Inches(w*0.4), Inches(0.06))
    ind.fill.solid(); ind.fill.fore_color.rgb = DGRAY; ind.line.fill.background()
    tab = s.shapes.add_shape(1, Inches(l+0.04), Inches(t+h-0.38),
                             Inches(w-0.08), Inches(0.26))
    tab.fill.solid(); tab.fill.fore_color.rgb = BG; tab.line.fill.background()
    for i, ic in enumerate(["🏠","🔍","🔔","👤"]):
        txt(s, ic, l+0.04+i*(w-0.08)/4, t+h-0.38, (w-0.08)/4, 0.26,
            7, align=PP_ALIGN.CENTER)
    return h

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)

# Büyük koyu mor blob - sağ üst
blob1 = s.shapes.add_shape(9, Inches(7.0), Inches(-2.0), Inches(8.5), Inches(8.5))
blob1.fill.solid(); blob1.fill.fore_color.rgb = RGBColor(0x18, 0x0A, 0x38)
blob1.line.fill.background()

# Koyu lacivert blob - sol alt
blob2 = s.shapes.add_shape(9, Inches(-2.0), Inches(4.0), Inches(6.0), Inches(6.0))
blob2.fill.solid(); blob2.fill.fore_color.rgb = RGBColor(0x0A, 0x10, 0x38)
blob2.line.fill.background()

# Sol dikey accent şerit
strip = s.shapes.add_shape(1, 0, 0, Inches(0.22), prs.slide_height)
strip.fill.solid(); strip.fill.fore_color.rgb = PURPLE; strip.line.fill.background()

# Uygulama ikonu
try:
    icon_path = r"C:\Users\EMRE\Desktop\concertly-mobile\mobile\assets\icon.png"
    s.shapes.add_picture(icon_path, Inches(1.2), Inches(1.5), Inches(1.4), Inches(1.4))
except:
    circle(s, 1.2, 1.5, 1.4, PURPLE)
    txt(s, "🎵", 1.2, 1.55, 1.4, 1.4, 36, align=PP_ALIGN.CENTER)

txt(s, "Concertly", 2.85, 1.48, 9, 1.2, 68, bold=True, color=WHITE)

# Gradient accent çizgi: mor → turuncu
ul = s.shapes.add_shape(1, Inches(2.85), Inches(2.78), Inches(6.5), Inches(0.07))
ul.fill.solid(); ul.line.fill.background()
apply_grad(ul, [(0, PURPLE), (1.0, ORANGE)], ang=0)

txt(s, "Konser & Etkinlik Keşif Sosyal Platformu",
    2.85, 2.95, 9.5, 0.55, 20, color=ORANGE)
txt(s, "Müziği birlikte keşfet, anları birlikte yaşa.",
    2.85, 3.58, 9.5, 0.5, 15, italic=True, color=LGRAY)

for i, (label, color) in enumerate([("React Native", PURPLE),
                                     ("Spring Boot",  BLUE),
                                     ("PostgreSQL",   ORANGE)]):
    pill(s, 2.85 + i*2.55, 4.55, 2.3, 0.4, color, label, 12)

phone_screenshot(s, SCREEN_HOME, 10.5, 0.8, 2.3)

txt(s, "Hazırlayan: Emre", 2.85, 6.9, 5, 0.42, 11, color=LGRAY)
txt(s, "2026", 11.8, 6.9, 1.3, 0.42, 11, color=LGRAY, align=PP_ALIGN.RIGHT)
pbar(s, 1, PURPLE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
circle(s, -1.0, -1.0, 4.0, RGBColor(0x18, 0x0A, 0x38))
circle(s, 11.0, 5.5,  3.0, RGBColor(0x0A, 0x10, 0x38))
slide_header(s, "Problem", "Neden Concertly'e ihtiyaç var?", PURPLE)

problems = [
    ("🗓️", "Dağınık Takvimler",
     "Konser bilgilerine ulaşmak için\nbir sürü farklı platform gerekiyor.", PURPLE),
    ("🎟️", "Yalnız Konserler",
     "Aynı müzik zevkini paylaşan\nbirini bulmak tesadüfe kalmış.", BLUE),
    ("📁", "Kayıp Anılar",
     "Gittiğin konserleri saklayacak\nbir koleksiyon sistemi yok.", ORANGE),
    ("🎯", "Kişiselleştirme Yok",
     "Müzik zevkine göre etkinlik\nönerisi sunan platform yok.", RED),
]

for i, (icon, title, desc, color) in enumerate(problems):
    l = 0.5 + i * 3.18
    card(s, l, 2.0, 3.0, 4.9, border_color=color)
    rect(s, l, 2.0, 3.0, 0.2, color)
    txt(s, icon,  l+1.1, 2.38, 1.0, 0.9,  34, align=PP_ALIGN.CENTER)
    txt(s, title, l+0.15, 3.38, 2.7, 0.5, 15, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, desc,  l+0.15, 3.95, 2.7, 1.4, 12, color=LGRAY, align=PP_ALIGN.CENTER)
    txt(s, "?",   l+1.25, 5.55, 0.5, 0.5, 28, bold=True, color=DGRAY, align=PP_ALIGN.CENTER)

pbar(s, 2, PURPLE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Çözüm
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
circle(s, 10.5, -1.0, 4.5, RGBColor(0x0A, 0x10, 0x38))
slide_header(s, "Çözüm", "Concertly — Hepsi Bir Arada Müzik & Etkinlik Platformu", BLUE)

features = [
    ("🎵", "Etkinlik Keşfi",       "Tür, şehir, tarih filtresi\nile Türkiye geneli etkinlikler", PURPLE),
    ("👥", "Sosyal Platform",      "Topluluklar, gönderi akışı\nve kullanıcı etkileşimi",         BLUE),
    ("🎸", "Konser Arkadaşı",      "Müzik uyumuna göre\nTinder tarzı eşleştirme",                ORANGE),
    ("🛂", "Konser Pasaportu",     "Katıldığın konserler,\nistatistikler ve rozetler",            RED),
    ("🗺️", "Etkinlik Haritası",    "İnteraktif harita ile\nkonum bazlı etkinlik keşfi",          PURPLE),
    ("🎧", "Spotify Entegrasyonu", "Dinleme geçmişine göre\nkişiselleştirilmiş öneriler",        BLUE),
]

positions = [(0.5,2.02),(4.6,2.02),(8.7,2.02),(0.5,4.57),(4.6,4.57),(8.7,4.57)]
for i, (icon, title, desc, color) in enumerate(features):
    l, t = positions[i]
    card(s, l, t, 3.9, 2.15, border_color=color)
    rect(s, l, t, 0.08, 2.15, color)
    txt(s, icon,  l+0.22, t+0.22, 0.65, 0.6, 24)
    txt(s, title, l+1.0,  t+0.2,  2.7,  0.4, 14, bold=True, color=color)
    txt(s, desc,  l+1.0,  t+0.65, 2.7,  1.3, 11, color=LGRAY)

pbar(s, 3, BLUE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Teknoloji Yığını
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Teknoloji Yığını", "Concertly'i oluşturan teknolojiler", ORANGE)

stack = [
    ("📱", "Mobil",       "React Native + Expo",   "iOS & Android cross-platform",  PURPLE),
    ("⚙️", "Backend",     "Spring Boot (Java 17)", "REST API, Security, JPA",       BLUE),
    ("🗄️", "Veritabanı",  "PostgreSQL",            "İlişkisel veri modeli",          ORANGE),
    ("🔐", "Kimlik Doğ.", "JWT",                   "Token tabanlı auth",             RED),
    ("🗺️", "Harita",      "React Native Maps",     "İnteraktif etkinlik haritası",  PURPLE),
    ("🎧", "Müzik",       "Spotify API",           "Sanatçı & öneri entegrasyonu",  BLUE),
    ("🎟️", "Etkinlik",    "Ticketmaster API",      "Etkinlik verisi import",        ORANGE),
]

for i, (icon, layer, tech, note, color) in enumerate(stack):
    t = 2.0 + i * 0.73
    card(s, 0.7, t, 12.0, 0.65)
    circle(s, 0.82, t+0.22, 0.22, color)
    txt(s, icon,  1.22, t+0.1,  0.55, 0.5, 18, align=PP_ALIGN.CENTER)
    txt(s, layer, 1.9,  t+0.13, 1.9,  0.4, 12, bold=True, color=LGRAY)
    txt(s, tech,  3.95, t+0.13, 3.8,  0.4, 13, bold=True, color=WHITE)
    txt(s, note,  7.9,  t+0.13, 4.6,  0.4, 11, color=LGRAY)
    pill(s, 12.35, t+0.16, 0.65, 0.3, color, "", 8)

pbar(s, 4, ORANGE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Kullanıcı Akışı
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
circle(s, -1.0, 4.5, 4.0, RGBColor(0x0A, 0x10, 0x38))
slide_header(s, "Kullanıcı Akışı", "Uygulamadaki kullanıcı yolculuğu", RED)

steps = [
    ("1", "Kayıt Ol",           "Kullanıcı adı,\ne-posta & şifre",          PURPLE),
    ("2", "Uygulama\nTanıtımı", "Özellikler\ntanıtılır",                     BLUE),
    ("3", "Ana Sayfa",          "Kişiselleştirilmiş\netkinlik akışı",         ORANGE),
    ("4", "Menü",               "Keşfet sekmesinden\ntüm özelliklere erişim", RED),
    ("5", "Özellikler",         "Etkinlikler, Harita,\nTopluluklar, Arkadaş…",PURPLE),
]

step_w = 2.3
for i, (num, title, desc, color) in enumerate(steps):
    l = 0.4 + i * (step_w + 0.12)
    t = 2.05
    card(s, l, t, step_w, 3.8, border_color=color)
    rect(s, l, t, step_w, 0.22, color)
    circle(s, l+step_w/2-0.3, t+0.35, 0.6, CARD_DK)
    circle(s, l+step_w/2-0.28, t+0.37, 0.56, color)
    txt(s, num, l+step_w/2-0.3, t+0.36, 0.6, 0.55, 18, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, title, l+0.1, t+1.05, step_w-0.2, 0.6, 14, bold=True, color=color,
        align=PP_ALIGN.CENTER)
    txt(s, desc,  l+0.1, t+1.72, step_w-0.2, 1.8, 11, color=LGRAY,
        align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        txt(s, "›", l+step_w+0.02, t+1.65, 0.16, 0.55, 24, bold=True,
            color=DGRAY, align=PP_ALIGN.CENTER)

txt(s, "★  Tür & sanatçı seçimi kayıt sırasında yapılır — tüm içerik kişiselleştirilir",
    0.7, 6.2, 12, 0.45, 12, italic=True, color=LGRAY)
pbar(s, 5, RED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Etkinlik Keşfi
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
slide_header(s, "Etkinlik Keşfi",
             "Türkiye genelinde konser & etkinlikleri keşfet", PURPLE)

phone_screenshot(s, SCREEN_EVENTS, 0.5, 1.75, 2.55)

features_disc = [
    ("🏠", "Ana Sayfa",       PURPLE,
     ["Şehir bazlı filtreleme (11 şehir)", "Kategori carousel",
      "Öne çıkan etkinlikler",             "Trend gönderiler"]),
    ("🎵", "Etkinlikler",     BLUE,
     ["Tür, şehir, tarih filtresi", "Sıralama seçenekleri",
      "Grid görünüm",                "Skeleton yükleme"]),
    ("📋", "Etkinlik Detayı", ORANGE,
     ["RSVP (Gidiyorum/İlgiliyorum)", "Takvime ekleme",
      "Arkadaşları gör",              "Yorum & beğeni"]),
    ("🗺️", "Harita",          RED,
     ["İnteraktif harita",  "Tür renk kodlaması",
      "Mesafe filtresi",    "Alt panel önizleme"]),
]

for i, (icon, title, color, bullets) in enumerate(features_disc):
    col = i % 2; row = i // 2
    l = 3.5 + col * 4.95
    t = 1.9 + row * 2.72
    card(s, l, t, 4.65, 2.52, border_color=color)
    rect(s, l, t, 4.65, 0.16, color)
    txt(s, f"{icon}  {title}", l+0.2, t+0.23, 4.2, 0.42, 14, bold=True, color=color)
    for j, b in enumerate(bullets):
        txt(s, f"▸  {b}", l+0.2, t+0.72+j*0.44, 4.25, 0.42, 11, color=WHITE)

pbar(s, 6, PURPLE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Sosyal Katman
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
circle(s, 10.5, 5.0, 4.0, RGBColor(0x0A, 0x10, 0x38))
slide_header(s, "Sosyal Katman", "Topluluk, etkileşim ve içerik paylaşımı", BLUE)

socials = [
    ("🔥", "Sosyal Akış",        PURPLE,
     ["Metin, fotoğraf, anket", "Trend / Takip edilen", "Beğeni & yorum"]),
    ("👥", "Topluluklar",        BLUE,
     ["Tür bazlı gruplar", "Üye ol / ayrıl", "Topluluk içi akış"]),
    ("🎤", "Sanatçı Profilleri", ORANGE,
     ["Bio & tür bilgisi", "Yaklaşan etkinlikler", "Takip et"]),
    ("🏛️", "Mekan Profilleri",   RED,
     ["Yıldız puanlama", "Kullanıcı yorumları", "Etkinlik takvimi"]),
    ("🔔", "Bildirimler",        PURPLE,
     ["Beğeni & yorum", "Takip bildirimleri", "Okunmamış sayaç"]),
    ("👤", "Kullanıcı Profili",  BLUE,
     ["Gönderi & etkinlik", "Takipçi/Takip", "Müzik profili"]),
]

phone_screenshot(s, SCREEN_FEED, 9.9, 1.75, 2.55)

positions6 = [(0.5,1.9),(3.42,1.9),(6.34,1.9),(0.5,4.6),(3.42,4.6),(6.34,4.6)]
for i, (icon, title, color, bullets) in enumerate(socials):
    l, t = positions6[i]
    card(s, l, t, 2.65, 2.4, border_color=color)
    rect(s, l, t, 2.65, 0.15, color)
    txt(s, f"{icon}  {title}", l+0.15, t+0.22, 2.35, 0.4, 13, bold=True, color=color)
    for j, b in enumerate(bullets):
        txt(s, f"• {b}", l+0.15, t+0.68+j*0.52, 2.35, 0.48, 11, color=WHITE)

pbar(s, 7, BLUE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Konser Pasaportu
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
circle(s, -1.0, -1.0, 5.0, RGBColor(0x18, 0x0A, 0x38))
circle(s, 10.5, 5.0,  4.0, RGBColor(0x0A, 0x10, 0x38))
slide_header(s, "Konser Pasaportu",
             "Gittiğin her konser pasaportunun bir sayfası", ORANGE)

card(s, 0.5, 2.0, 7.5, 5.1, border_color=ORANGE)
rect(s, 0.5, 2.0, 7.5, 0.22, ORANGE)
txt(s, "🛂", 0.75, 2.32, 1.0, 0.9, 36)
txt(s, "Nedir?", 1.85, 2.37, 6.0, 0.4, 16, bold=True, color=ORANGE)
txt(s, "Konserlere katıldıkça büyüyen kişisel koleksiyonun.",
    1.85, 2.82, 6.0, 0.4, 13, italic=True, color=LGRAY)

passport_feats = [
    ("📅", "Yıllara Göre",  "Konserler yıl başlıkları altında\ngruplandırılır"),
    ("📊", "İstatistikler", "Toplam konser, sanatçı\nve şehir sayısı"),
    ("🏅", "Rozetler",      "Doğrulanmış katılım\nve başarı rozetleri"),
    ("🔗", "Paylaşım",      "Pasaportunun özetini\nsosyal medyada paylaş"),
]

for i, (icon, title, desc) in enumerate(passport_feats):
    col = i % 2; row = i // 2
    l = 0.75 + col * 3.75
    t = 3.47 + row * 1.45
    card2(s, l, t, 3.5, 1.3, border_color=ORANGE)
    txt(s, icon,  l+0.15, t+0.15, 0.6, 0.55, 22)
    txt(s, title, l+0.85, t+0.12, 2.5, 0.35, 13, bold=True, color=ORANGE)
    txt(s, desc,  l+0.85, t+0.52, 2.5, 0.7,  11, color=LGRAY)

phone_screenshot(s, SCREEN_PASSPORT, 8.55, 1.75, 2.55)

stats = [("21","Konser",ORANGE),("14","Sanatçı",PURPLE),("6","Şehir",BLUE)]
for i, (num, label, color) in enumerate(stats):
    l = 11.4 + i * 0.65
    card2(s, l, 2.2, 0.6, 1.1)
    txt(s, num,   l, 2.28, 0.6, 0.42, 20, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, label, l, 2.72, 0.6, 0.3,  8,  color=LGRAY, align=PP_ALIGN.CENTER)

pbar(s, 8, ORANGE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Konser Arkadaşı
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
circle(s, 10.0, -1.5, 5.5, RGBColor(0x18, 0x0A, 0x38))
slide_header(s, "Konser Arkadaşı Eşleştirme",
             "Müzik zevkine göre Tinder tarzı eşleştirme", RED)

phone_screenshot(s, SCREEN_BUDDY, 0.5, 1.75, 2.55)

card(s, 3.5, 1.9, 9.4, 5.2, border_color=RED)
rect(s, 3.5, 1.9, 9.4, 0.22, RED)

buddy_feats = [
    ("🃏", "Kart Sistemi",       RED,
     "Kaydırılabilir kullanıcı kartları\nfotoğraf, bio ve müzik bilgisiyle"),
    ("🎯", "Uyum Skoru",         BLUE,
     "0–100 arası tür eşleşme skoru\n'Harika / İyi / Farklı Zevk' etiketleri"),
    ("💫", "Karşılıklı Eşleşme", ORANGE,
     "İki taraf da beğenince\nkonfeti animasyonu ve 'Eşleşme!' ekranı"),
    ("👁️", "Profil Keşfi",        PURPLE,
     "Eşleştiğin kişinin\nprofiline giderek daha fazla tanı"),
]

for i, (icon, title, color, desc) in enumerate(buddy_feats):
    col = i % 2; row = i // 2
    l = 3.7 + col * 4.6
    t = 2.18 + row * 2.45
    card2(s, l, t, 4.35, 2.2, border_color=color)
    txt(s, icon,  l+0.2,  t+0.18, 0.7,  0.6,  26)
    txt(s, title, l+1.05, t+0.15, 3.1,  0.42, 14, bold=True, color=color)
    txt(s, desc,  l+1.05, t+0.62, 3.1,  1.4,  12, color=LGRAY)

pbar(s, 9, RED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Admin Paneli
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
circle(s, -0.5, 5.0, 4.0, RGBColor(0x0A, 0x10, 0x38))
slide_header(s, "Yönetim Paneli",
             "Platform yönetimi için tam kapsamlı admin altyapısı", BLUE)

stats_adm = [
    ("👥","Kullanıcılar",PURPLE),("🎵","Etkinlikler",BLUE),
    ("📝","Gönderiler", ORANGE), ("🎟️","Katılımlar",  RED),
]
for i, (icon, label, color) in enumerate(stats_adm):
    l = 0.5 + i * 3.2
    card(s, l, 2.0, 3.0, 1.1, border_color=color)
    txt(s, icon,  l+0.2, 2.05, 0.6, 0.55, 24)
    txt(s, label, l+0.9, 2.1,  2.0, 0.4,  14, bold=True, color=color)
    txt(s, "Gerçek zamanlı", l+0.9, 2.54, 2.0, 0.3, 9, color=LGRAY)

admin_cards = [
    ("🎵 Etkinlik Yönetimi", PURPLE,
     ["Etkinlik oluştur, düzenle, sil",
      "Bekleyen etkinlikleri onayla/reddet",
      "Sanatçı & mekan ataması"]),
    ("👥 Kullanıcı Yönetimi", BLUE,
     ["Kullanıcıları listele & ara",
      "Hesap ban / aktifleştir",
      "Admin rolü ata / kaldır"]),
    ("📝 İçerik Moderasyonu", ORANGE,
     ["Tüm gönderileri görüntüle",
      "Şikayet edilen içerikleri incele",
      "Gönderi sil & geri yükle"]),
    ("📊 İstatistik Paneli", RED,
     ["Toplam/aktif/yeni kullanıcı",
      "Bekleyen etkinlik sayısı",
      "Katılım & etkileşim metrikleri"]),
]

for i, (title, color, bullets) in enumerate(admin_cards):
    col = i % 2; row = i // 2
    l = 0.5 + col * 6.5
    t = 3.38 + row * 2.0
    card(s, l, t, 6.1, 1.82, border_color=color)
    rect(s, l, t, 6.1, 0.15, color)
    txt(s, title, l+0.2, t+0.22, 5.7, 0.4, 14, bold=True, color=color)
    for j, b in enumerate(bullets):
        txt(s, f"▸  {b}", l+0.2, t+0.68+j*0.36, 5.7, 0.33, 11, color=WHITE)

pbar(s, 10, BLUE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Durum & Yol Haritası
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
circle(s, 10.5, -0.5, 4.5, RGBColor(0x0A, 0x10, 0x38))
slide_header(s, "Mevcut Durum & Yol Haritası",
             "Neredeyiz, nereye gidiyoruz?", PURPLE)

card(s, 0.5, 1.95, 6.0, 5.2, border_color=PURPLE)
rect(s, 0.5, 1.95, 6.0, 0.22, PURPLE)
txt(s, "✅  Tamamlanan", 0.72, 2.22, 5.5, 0.45, 17, bold=True, color=PURPLE)

done = [
    "Tüm ana ekranlar ve özellikler",
    "Konser Pasaportu & Arkadaş Eşleştirme",
    "Sosyal akış, topluluklar, profiller",
    "Spotify & Ticketmaster entegrasyonu",
    "Admin yönetim paneli",
    "TR / EN çoklu dil desteği",
    "Karanlık & Aydınlık tema",
    "Skeleton loader & Lottie animasyonları",
]
for i, item in enumerate(done):
    txt(s, f"✓  {item}", 0.72, 2.75+i*0.53, 5.5, 0.48, 12, color=WHITE)

card(s, 7.05, 1.95, 6.0, 5.2, border_color=ORANGE)
rect(s, 7.05, 1.95, 6.0, 0.22, ORANGE)
txt(s, "🔜  Yakında", 7.27, 2.22, 5.5, 0.45, 17, bold=True, color=ORANGE)

coming = [
    "Canlı Sohbet Odaları (WebSocket)",
    "Bilet Uyarıları & Anlık Bildirimler",
    "App Store & Google Play yayını",
]
for i, item in enumerate(coming):
    txt(s, f"○  {item}", 7.27, 2.75+i*0.65, 5.5, 0.55, 13, color=WHITE)

pbar(s, 11, PURPLE)

# ── Galeri slayt yardımcısı ───────────────────────────────────────────────
def gallery_slide(title, subtitle, screens, slide_num, accent=PURPLE):
    s = prs.slides.add_slide(BLANK)
    slide_bg(s)
    slide_header(s, title, subtitle, accent)
    ph = 4.7
    pw = ph / 2.167
    gap = 0.4
    total_w = len(screens) * pw + (len(screens) - 1) * gap
    x0 = (W - total_w) / 2
    y0 = 2.05
    for i, (path, label) in enumerate(screens):
        lx = x0 + i * (pw + gap)
        sd = s.shapes.add_shape(1, Inches(lx+0.06), Inches(y0+0.08),
                                 Inches(pw+0.04), Inches(ph+0.04))
        sd.fill.solid(); sd.fill.fore_color.rgb = CARD_DK; sd.line.fill.background()
        frame = s.shapes.add_shape(5, Inches(lx-0.05), Inches(y0-0.05),
                                    Inches(pw+0.1), Inches(ph+0.1))
        frame.fill.solid(); frame.fill.fore_color.rgb = RGBColor(0x1A, 0x18, 0x30)
        frame.line.color.rgb = DGRAY; frame.line.width = Inches(0.012)
        try:
            s.shapes.add_picture(path, Inches(lx), Inches(y0), Inches(pw), Inches(ph))
        except Exception:
            fb = s.shapes.add_shape(1, Inches(lx), Inches(y0), Inches(pw), Inches(ph))
            fb.fill.solid(); fb.fill.fore_color.rgb = CARD2; fb.line.fill.background()
        txt(s, label, lx, y0 + ph + 0.06, pw, 0.25, 10,
            bold=True, color=accent, align=PP_ALIGN.CENTER)
    pbar(s, slide_num, accent)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Galeri: Giriş & Kayıt
# ═══════════════════════════════════════════════════════════════════════════
gallery_slide("Giriş & Kayıt",
              "Kullanıcı kimlik doğrulama ekranları",
              [(SCREEN_LOGIN,    "Giriş Yap"),
               (SCREEN_REGISTER, "Kayıt Ol"),
               (SCREEN_ADMIN,    "Admin Paneli")],
              12, RED)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Galeri: Ana Ekranlar
# ═══════════════════════════════════════════════════════════════════════════
gallery_slide("Ana Ekranlar",
              "Etkinlik keşfi ve sosyal akış",
              [(SCREEN_HOME,   "Ana Sayfa"),
               (SCREEN_EVENTS, "Etkinlikler"),
               (SCREEN_FEED,   "Sosyal Feed")],
              13, PURPLE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Galeri: Topluluk & Sosyal
# ═══════════════════════════════════════════════════════════════════════════
gallery_slide("Topluluk & Sosyal",
              "Topluluklar, profil ve rozetler",
              [(SCREEN_COMMUNITIES, "Topluluklar"),
               (SCREEN_PROFILE,     "Profil"),
               (SCREEN_BADGES,      "Rozetler")],
              14, BLUE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Galeri: Özel Özellikler
# ═══════════════════════════════════════════════════════════════════════════
gallery_slide("Özel Özellikler",
              "Konser arkadaşı, pasaport ve müzik profili",
              [(SCREEN_BUDDY,    "Konser Arkadaşı"),
               (SCREEN_PASSPORT, "Konser Pasaportu"),
               (SCREEN_MUSIC,    "Müzik Profili")],
              15, ORANGE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Galeri: Konser Haritası
# ═══════════════════════════════════════════════════════════════════════════
gallery_slide("Konser Haritası",
              "İnteraktif harita ile konum bazlı etkinlik keşfi",
              [(SCREEN_MAP_TR,  "Türkiye Geneli"),
               (SCREEN_MAP_ANK, "Ankara Yakını")],
              16, PURPLE)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Kapanış
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)

blob_a = s.shapes.add_shape(9, Inches(-2.0), Inches(-2.0), Inches(7), Inches(7))
blob_a.fill.solid(); blob_a.fill.fore_color.rgb = RGBColor(0x18, 0x0A, 0x38)
blob_a.line.fill.background()

blob_b = s.shapes.add_shape(9, Inches(9.0), Inches(3.5), Inches(7), Inches(7))
blob_b.fill.solid(); blob_b.fill.fore_color.rgb = RGBColor(0x0A, 0x10, 0x38)
blob_b.line.fill.background()

strip = s.shapes.add_shape(1, 0, 0, Inches(0.22), prs.slide_height)
strip.fill.solid(); strip.fill.fore_color.rgb = PURPLE; strip.line.fill.background()

try:
    icon_path = r"C:\Users\EMRE\Desktop\concertly-mobile\mobile\assets\icon.png"
    s.shapes.add_picture(icon_path, Inches(5.67), Inches(1.3), Inches(2.0), Inches(2.0))
except:
    circle(s, 5.67, 1.3, 2.0, PURPLE)

txt(s, "Teşekkürler", 1.5, 3.5, 10.33, 1.1, 62, bold=True, color=WHITE,
    align=PP_ALIGN.CENTER)

ul2 = s.shapes.add_shape(1, Inches(4.0), Inches(4.72), Inches(5.33), Inches(0.07))
ul2.fill.solid(); ul2.line.fill.background()
apply_grad(ul2, [(0, PURPLE), (1.0, ORANGE)], ang=0)

txt(s, "Müziği birlikte keşfet, anları birlikte yaşa.",
    1.5, 4.88, 10.33, 0.55, 18, italic=True, color=ORANGE, align=PP_ALIGN.CENTER)
txt(s, "Sorularınız?", 1.5, 5.5, 10.33, 0.5, 16, color=LGRAY, align=PP_ALIGN.CENTER)

for i, (label, color) in enumerate([("React Native", PURPLE),
                                     ("Spring Boot",  BLUE),
                                     ("PostgreSQL",   ORANGE),
                                     ("Spotify API",  RED)]):
    total_w = 4 * 2.3 + 3 * 0.15
    start = (W - total_w) / 2
    pill(s, start + i * 2.45, 6.52, 2.2, 0.36, color, label, 11)

pbar(s, 17, PURPLE)

# ── Kaydet ───────────────────────────────────────────────────────────────────
output = r"C:\Users\EMRE\Desktop\Concertly_Sunum_v7.pptx"
prs.save(output)
print(f"Kaydedildi: {output}")
